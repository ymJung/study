""" requirements.txt
# PDF 처리 및 PPT 처리 관련
PyMuPDF==1.25.3
python-pptx==1.0.2
# 이미지 처리 및 OCR
opencv-python==4.11.0.86
numpy==1.26.4
Pillow==11.0.0
pytesseract==0.3.13

# LangChain 및 벡터 DB 관련 (사용하는 LangChain 모듈에 따라 버전 조정)
langchain==0.3.17
langchain-community==0.3.16
langchain-openai==0.3.4
langchain-huggingface==0.1.2
faiss-cpu==1.10.0
"""
""" Dockerfile
FROM python:3.10-slim

WORKDIR /app

RUN apt-get update && \
    apt-get install -y tesseract-ocr tesseract-ocr-kor libgl1-mesa-glx libglib2.0-0 && \
    apt-get clean && rm -rf /var/lib/apt/lists/*

COPY requirements.txt .
RUN pip install --no-cache-dir -r requirements.txt
COPY . .
CMD ["python", "llm_rag.py"]

"""


import os
import io
import json
import fitz  # pip install PyMuPDF
from pptx import Presentation  # pip install python-pptx
import cv2
import numpy as np
from PIL import Image
import pytesseract  # sudo apt-get install tesseract-ocr, pip install pytesseract
from langchain.text_splitter import RecursiveCharacterTextSplitter 
from langchain.docstore.document import Document
from langchain_huggingface import HuggingFaceEmbeddings
from langchain_community.vectorstores import FAISS  # pip install faiss-cpu
from langchain.vectorstores import FAISS
from langchain_openai import ChatOpenAI
from langchain.chains import RetrievalQA
import configparser
import re

# ============================================================
# 설정 파일 로드 (config.cfg 파일 필요: [openai] TOKEN = <API_KEY>)
# ============================================================
config = configparser.ConfigParser()
config.read('config.cfg')
openai_api_key = config['openai']['TOKEN']

# ============================================================
# 전처리 함수: 이미지 OCR 전처리 및 한글 띄어쓰기 정리
# ============================================================
def preprocess_image_for_ocr_alternative(image):
    image_cv = cv2.cvtColor(np.array(image), cv2.COLOR_RGB2BGR)
    gray = cv2.cvtColor(image_cv, cv2.COLOR_BGR2GRAY)
    clahe = cv2.createCLAHE(clipLimit=2.0, tileGridSize=(8,8))
    equalized = clahe.apply(gray)
    sharpening_kernel = np.array([[0, -1, 0],
                                  [-1, 5, -1],
                                  [0, -1, 0]])
    sharpened = cv2.filter2D(equalized, -1, sharpening_kernel)
    processed_image = Image.fromarray(sharpened)
    return processed_image

def clean_hangul_spacing(text): # 한글 문자 사이에 있는 불필요한 공백을 제거
    return re.sub(r'(?<=[가-힣])\s+(?=[가-힣])', '', text)

# ============================================================
# 문서 전처리 및 청킹 함수 (PDF PPT)
# PDF 파일에서 텍스트 추출(텍스트 레이어 또는 OCR 적용)
# 라인별/컬럼별 그룹화와 후처리를 통해 Document  생성
# 메타데이터에 페이지 번호와 에이전트 정보 등을 기록
# ============================================================
def process_pdf(file_path, agent="general"):
    docs = []
    splitter = RecursiveCharacterTextSplitter(chunk_size=500, chunk_overlap=100)
    try:
        with fitz.open(file_path) as pdf:
            for page in pdf:
                text = page.get_text("text").strip()
                if not text:
                    # OCR 적용: 텍스트 레이어가 없는 경우
                    try:
                        pix = page.get_pixmap()
                        img = Image.frombytes("RGB", (pix.width, pix.height), pix.samples)
                        processed_img = preprocess_image_for_ocr_alternative(img)
                        custom_config = r'--oem 3 --psm 6'
                        ocr_data = pytesseract.image_to_data(
                            processed_img,
                            lang="eng+kor",
                            config=custom_config,
                            output_type=pytesseract.Output.DICT
                        )
                        # 1. 각 라인별로 단어 정보 수집 (신뢰도 60 이상)
                        horizontal_gap_threshold = 50  # 단어 사이 가로 간격 임계값
                        lines_dict = {}
                        num_boxes = len(ocr_data['text'])
                        for i in range(num_boxes):
                            try:
                                conf = float(ocr_data['conf'][i])
                            except Exception:
                                conf = 0.0
                            if conf > 60 and ocr_data['text'][i].strip():
                                line_num = ocr_data['line_num'][i]
                                if line_num not in lines_dict:
                                    lines_dict[line_num] = []
                                left = ocr_data['left'][i]
                                width = ocr_data['width'][i]
                                word = ocr_data['text'][i].strip()
                                lines_dict[line_num].append((left, width, word))
                        # 2. 각 라인 내 단어들을 좌측 정렬 후, 가로 간격에 따라 그룹 분리
                        grouped_lines = []
                        for line in sorted(lines_dict.keys()):
                            words = sorted(lines_dict[line], key=lambda x: x[0])
                            groups = []
                            current_group = []
                            for word_info in words:
                                if not current_group:
                                    current_group.append(word_info)
                                else:
                                    last_word = current_group[-1]
                                    gap = word_info[0] - (last_word[0] + last_word[1])
                                    if gap > horizontal_gap_threshold:
                                        groups.append(" ".join(w for _, _, w in current_group))
                                        current_group = [word_info]
                                    else:
                                        current_group.append(word_info)
                            if current_group:
                                groups.append(" ".join(w for _, _, w in current_group))
                            grouped_lines.append(groups)
                        # 3. 각 라인의 컬럼별 그룹 병합: 각 컬럼(인덱스)별로 모든 라인의 텍스트 합침
                        max_cols = max(len(groups) for groups in grouped_lines)
                        merged_regions = []
                        for col in range(max_cols):
                            region_texts = []
                            for groups in grouped_lines:
                                if col < len(groups):
                                    region_texts.append(groups[col])
                            merged_regions.append(" ".join(region_texts))
                        # 4. 모든 컬럼을 하나의 텍스트로 합침 및 후처리
                        page_text = " ".join(merged_regions)
                        page_text = clean_hangul_spacing(page_text)
                        docs.append(Document(
                            page_content=page_text,
                            metadata={
                                "source": os.path.basename(file_path),
                                "page": page.number + 1,
                                "agent": agent,
                                "doc_type": "PDF"
                            }
                        ))
                    except Exception as ocr_err:
                        print(f"[{os.path.basename(file_path)}] OCR 실패 (페이지 {page.number + 1}): {ocr_err}")
                        text = ""
                else:
                    # 텍스트 레이어가 있는 경우: 기존 텍스트를 청킹
                    chunks = splitter.split_text(text)
                    for chunk in chunks:
                        docs.append(Document(
                            page_content=chunk,
                            metadata={
                                "source": os.path.basename(file_path),
                                "page": page.number + 1,
                                "agent": agent,
                                "doc_type": "PDF"
                            }
                        ))
    except Exception as e:
        print(f"PDF 처리 중 오류 발생 ({file_path}): {e}")
    return docs

def extract_text_from_shape(shape): 
    text = ""
    if hasattr(shape, "text") and shape.text:
        text += shape.text + "\n"
    if hasattr(shape, "shapes"):
        for subshape in shape.shapes:
            text += extract_text_from_shape(subshape) # 필요하면 재귀귀
    return text

def process_ppt(file_path, agent="general"):
    docs = []
    splitter = RecursiveCharacterTextSplitter(chunk_size=500, chunk_overlap=100)
    try:
        prs = Presentation(file_path)
        for i, slide in enumerate(prs.slides):
            slide_text = ""            
            for shape in slide.shapes:                                
                extracted = extract_text_from_shape(shape) # 전부추출                
                if extracted.strip():
                    slide_text += extracted + "\n"
                # 이미지 OCR 처리
                if shape.shape_type == 13:
                    try:
                        image_blob = shape.image.blob
                        img = Image.open(io.BytesIO(image_blob))
                        processed_img = preprocess_image_for_ocr_alternative(img)
                        custom_config = r'--oem 3 --psm 6'
                        ocr_text = pytesseract.image_to_string(processed_img, lang="eng+kor", config=custom_config)
                        if ocr_text.strip():
                            slide_text += ocr_text + "\n"
                    except Exception as img_err:
                        print(f"[{os.path.basename(file_path)}] PPT 이미지 OCR 오류 (슬라이드 {i+1}): {img_err}")
            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                for shape in notes_slide.shapes:
                    slide_text += extract_text_from_shape(shape) + "\n"
                    
            if slide_text.strip():
                slide_text = clean_hangul_spacing(slide_text)
                chunks = splitter.split_text(slide_text)
                for chunk in chunks:
                    docs.append(Document(
                        page_content=chunk,
                        metadata={
                            "source": os.path.basename(file_path),
                            "slide": i + 1,
                            "agent": agent,
                            "doc_type": "PPT"
                        }
                    ))
    except Exception as e:
        print(f"PPT 처리 중 오류 발생 ({file_path}): {e}")
    return docs

# ============================================================
# 임베딩 생성 및 벡터 DB 구축
# 여러 Document 리스트를, FAISS 벡터 스토어를 생성
# ============================================================
def setup_vector_store(documents):
    all_docs = []
    for doc_list in documents:
        all_docs.extend(doc_list) # 합쳐서 처리리
    embeddings = HuggingFaceEmbeddings(model_name="sentence-transformers/paraphrase-multilingual-mpnet-base-v2") # 다국어 지원모델
    vectorstore = FAISS.from_documents(all_docs, embeddings)
    return vectorstore

# ============================================================
# RAG 체인 생성 (벡터 스토어 + LLM 연결)
# LLM 답변과 함께 관련 근거(출처)를 반환
# ============================================================
def create_retrieval_qa(vectorstore, llm=None):
    if llm is None:
        llm = ChatOpenAI(model_name="gpt-4", temperature=0, openai_api_key=openai_api_key)
    qa = RetrievalQA.from_chain_type(
        llm=llm,
        chain_type="stuff",
        retriever=vectorstore.as_retriever(search_kwargs={"k": 3}),
        return_source_documents=True
    )
    return qa

# ============================================================
# 일반 질문 응답 (Retrieval 없이 단순 LLM 활용)
# ============================================================
def create_general_qa(llm=None):
    if llm is None:
        llm = ChatOpenAI(model_name="gpt-4", temperature=0, openai_api_key=openai_api_key)
    def answer(query):
        response = llm.invoke(query)
        return {"result": response, "source_documents": []}
    return answer

# ============================================================
# 에이전트 라우팅 규칙 로드 (외부 JSON 파일 지원)
# ============================================================
def load_agent_routing_rules(filepath="agent_routing_rules.json"):
    """
    에이전트 라우팅 규칙을 외부 JSON 파일로부터 로드합니다.
    파일이 없거나 오류가 발생하면 기본 
    """
    try:
        with open(filepath, "r", encoding="utf-8") as f:
            rules = json.load(f)
            return rules
    except Exception as e:
        print(f"에이전트 라우팅 규칙 파일 로드 실패({filepath}): {e}")
        return [
            {"agent": "deepseek", "keywords": ["강화학습", "deepseek"]},
            {"agent": "brochure", "keywords": ["서비스"]},
            {"agent": "rule", "keywords": ["회사"]}
        ]

# 전역 변수로 라우팅 규칙 로드
AGENT_ROUTING_RULES = load_agent_routing_rules()

def route_agent(query):
    """
    쿼리 문자열을 받아서, 미리 정의한 키워드 기반 규칙(외부 파일 또는 기본값)에 따라 에이전트를 선택합니다.
    우선순위는 AGENT_ROUTING_RULES 리스트의 순서를 따릅니다.
    """
    lower_query = query.lower()
    for rule in AGENT_ROUTING_RULES:
        for keyword in rule["keywords"]:
            if keyword.lower() in lower_query:
                return rule["agent"]
    return "general"

# ============================================================
# 사용자 질의응답 루프 (엔터 무시, 'exit' 입력 시 종료)
# ============================================================
def query_loop(qa_deepseek=None, qa_brochure=None, qa_rule=None, general_qa=None):
    print("문서 기반 질의응답 시스템을 시작합니다. (종료하려면 'exit' 입력)")
    while True:
        query = input("\n🔹 질문을 입력하세요: ").strip()
        if query == "":
            continue  # 빈 입력은 무시
        if query.lower() == "exit":
            print("프로그램을 종료합니다.")
            break

        # 에이전트 선택: 외부 규칙에 따라 deepseek, brochure, rule, general 중 선택
        agent = route_agent(query)
        if agent == "deepseek":
            qa_chain = qa_deepseek
        elif agent == "brochure":
            qa_chain = qa_brochure
        elif agent == "rule":
            qa_chain = qa_rule
        else:
            qa_chain = None  # 일반 agent

        if qa_chain:
            result = qa_chain.invoke({"query": query})
            answer = result.get("result", "")
            source_docs = result.get("source_documents", [])
            if not source_docs:
                answer = "해당 질문에 대한 답변을 찾을 수 없습니다."
                source_str = "General Agent"
            else:
                sources = set()
                for doc in source_docs:
                    meta = doc.metadata
                    if "page" in meta:
                        sources.add(f"{meta.get('source', '')} {meta.get('page', '')}페이지")
                    elif "slide" in meta:
                        sources.add(f"{meta.get('source', '')} {meta.get('slide', '')}슬라이드")
                source_str = ", ".join(sources)
        else:
            # 일반 agent: retrieval 없이 LLM 단독 사용
            res = general_qa(query)
            answer = res.get("result", "")
            source_str = "General Agent (No Document Source)"
        
        print("답변:", answer)
        print("출처:", source_str)

# ============================================================
# main: 문서 전처리, 벡터 스토어 구축, RAG 체인 생성 및 질의응답 루프 실행
# ============================================================
def main():
    # 문서 처리: 각 에이전트별로 메타데이터에 "agent" 기록
    docs_deepseek = process_pdf("pdf_file1.pdf", agent="deepseek")
    docs_brochure = process_pdf("pdf_file2.pdf", agent="brochure")
    docs_rule = process_ppt("ppt_file1.pptx", agent="rule")
    # 벡터 DB 구성: 각 에이전트별 별도 벡터 스토어 구축
    vectorstore_deepseek = setup_vector_store([docs_deepseek])
    vectorstore_brochure = setup_vector_store([docs_brochure])
    vectorstore_rule = setup_vector_store([docs_rule])
    
    # RAG 체인 생성
    qa_deepseek = create_retrieval_qa(vectorstore_deepseek)
    qa_brochure = create_retrieval_qa(vectorstore_brochure)
    qa_rule = create_retrieval_qa(vectorstore_rule)
    # 일반 에이전트 (retrieval 미적용)
    general_qa = create_general_qa()
    
    # 사용자 질의응답 루프 실행
    query_loop(qa_deepseek=qa_deepseek, qa_brochure=qa_brochure, qa_rule=qa_rule, general_qa=general_qa)


if __name__ == '__main__':
    main()
